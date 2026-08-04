"""Пайплайн генерации презентации на настоящем PostgreSQL.

Ollama и ChromaDB не поднимаются: подменены ретривал и фабрика ModelManager
(build_presentation_model_manager) — ровно две внешние зависимости пайплайна.
Подменяется именно ФАБРИКА, а не класс: клиента пути презентаций строит она, и
её же настоящую проверяет сторож связи потолков (tests/
test_presentation_call_timeout_single_source.py).

Всё остальное настоящее, включая рендер
python-pptx, запись файла и строку в журнале блокнота: именно на стыке «файл
на диске — строка в базе» и живут интересные ошибки.

Что проверяется:

* источники перечитываются в момент генерации, а не при постановке в очередь
  (presentation.no_sources, когда их удалили, пока задача стояла);
* файл появляется атомарно и ДО commit'а со status='ready' — строка не имеет
  права обещать файл, которого нет;
* временный файл подчищается при любом исходе;
* коды отказов: generation_failed, ollama_unavailable, generation_timeout —
  и то, что снятая таймаутом джоба не оставляет за собой файлов;
* таймаут живёт на уровне ВЫЗОВА, а не джобы: зависший вызов снимается своим
  бюджетом, повторная попытка имеет собственный, и error_text называет стадию
  и номер слайда. Подменённая модель не ходит по HTTP, поэтому здесь работает
  ВТОРОЙ эшелон — wait_for на LLM_CALL_WATCHDOG_TIMEOUT; в бою первым сдаётся
  клиент со своим LLM_CALL_TIMEOUT, и это ровно тот случай, ради которого
  страховка и оставлена: зависание, которого клиент не видит;
* дайджест уже написанного доезжает до второго слайд-вызова (мера 2 правила
  «не добивать»).
"""

import asyncio
import json
import os
import sys
import tempfile
import unittest
from contextlib import asynccontextmanager
from time import perf_counter
from unittest.mock import MagicMock, patch

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from dbfixtures import DatabaseBackedTestCase  # noqa: E402

from app.core.exceptions import ExternalServiceError, PresentationErrors  # noqa: E402
from app.modules.presentations.constants import (  # noqa: E402
    SLIDE_COUNT_MIN,
    STATUS_ERROR,
    STATUS_GENERATING,
    STATUS_QUEUED,
    STATUS_READY,
    presentation_job_timeout,
)
from app.modules.presentations.llm_schemas import content_section_count  # noqa: E402
from app.modules.presentations import service as presentation_service  # noqa: E402
from app.modules.presentations.worker import PresentationWorker  # noqa: E402
from app.shared.models import (  # noqa: E402
    Document,
    Log,
    Notebook,
    Presentation,
    User,
)
from app.shared.settings.config import settings as app_settings  # noqa: E402

EMBEDDING_MODEL = "qwen3-embedding:8b"

# Заказ ровно на нижней границе допустимого (SLIDE_COUNT_MIN): самая короткая
# колода, которую вообще можно заказать через API, — значит и самый дешёвый
# прогон полного пайплайна. Число секций в плане отсюда же выводится, а не
# выписано: разъедься они, план не прошёл бы валидацию, и тест ловил бы это
# как «модель вернула не то».
DECK_SLIDES = SLIDE_COUNT_MIN

PLAN_SECTIONS = [
    {"heading": "Кто имеет право", "search_query": "право на льготу"},
    {"heading": "Как оформить", "search_query": "порядок оформления"},
    {"heading": "Куда обращаться", "search_query": "куда обращаться"},
]
assert len(PLAN_SECTIONS) == content_section_count(DECK_SLIDES), (
    "план фикстуры разошёлся с числом слайдов заказа"
)

PLAN_JSON = json.dumps(
    {"title": "Налоговые льготы", "sections": PLAN_SECTIONS},
    ensure_ascii=False,
)


def slide_json(heading: str, bullets: list[str], chunk_id: int) -> str:
    return json.dumps(
        {
            "heading": heading,
            "bullets": bullets,
            "citations": [{"source_id": 1, "chunk_id": chunk_id}],
        },
        ensure_ascii=False,
    )


class FakeModelManager:
    """Очередь заранее заготовленных ответов вместо Ollama.

    Записывает полученные сообщения: по ним проверяется, что в слайд-вызов
    действительно уехал дайджест уже написанного.
    """

    def __init__(self, responses: list, calls: list) -> None:
        self._responses = responses
        self.calls = calls

    async def chat(self, *, model=None, messages=None, num_ctx=None) -> str:
        self.calls.append(messages)
        if not self._responses:
            raise AssertionError("модель вызвали больше раз, чем ожидалось")
        response = self._responses.pop(0)
        if isinstance(response, Exception):
            raise response
        return response


class PresentationPipelineTestCase(DatabaseBackedTestCase):
    async def asyncSetUp(self) -> None:
        await super().asyncSetUp()

        self.user: User = await self.make_user("owner", "user")
        self.as_user(self.user)
        self.notebook: Notebook = await self.seed(
            Notebook(name="Налоги", domain_profile="tax", owner_id=self.user.id)
        )
        self.document: Document = await self.seed(
            Document(
                name="Кодекс.pdf",
                path=self.make_file("Кодекс.pdf"),
                size=10,
                status="indexed",
                notebook_id=self.notebook.id,
                owner_id=self.user.id,
            )
        )

        env_patcher = patch.object(
            app_settings, "OLLAMA_MODEL_EMBEDDING", EMBEDDING_MODEL
        )
        env_patcher.start()
        self.addCleanup(env_patcher.stop)

        for target in (
            "app.modules.presentations.worker.session_context",
            "app.modules.presentations.service.session_context",
        ):
            patcher = patch(target, self._worker_session)
            patcher.start()
            self.addCleanup(patcher.stop)

        # Каталог файлов — временный: настоящий data/presentations тесты не
        # трогают.
        self.storage = tempfile.TemporaryDirectory()
        self.addCleanup(self.storage.cleanup)
        storage_patcher = patch.object(
            presentation_service, "PRESENTATION_STORAGE_DIR", self.storage.name
        )
        storage_patcher.start()
        self.addCleanup(storage_patcher.stop)

        rag_patcher = patch.object(presentation_service, "RAGService")
        rag_patcher.start()
        self.addCleanup(rag_patcher.stop)

        settings_patcher = patch.object(
            presentation_service, "RuntimeSettingsService", MagicMock()
        )
        runtime_settings = settings_patcher.start()
        runtime_settings.get_settings.return_value = {"chat_model": "test-model"}
        self.addCleanup(settings_patcher.stop)

        self.model_calls: list = []
        self.retrieval_queries: list[str] = []
        self.worker = PresentationWorker(poll_interval=0.05)
        self.addAsyncCleanup(self.worker.stop)

    @asynccontextmanager
    async def _worker_session(self):
        async with self.session_factory() as session:
            try:
                yield session
            except Exception:
                await session.rollback()
                raise
            finally:
                await session.close()

    # --- подмены ---

    def use_model(self, responses: list) -> None:
        manager = FakeModelManager(responses, self.model_calls)
        patcher = patch.object(
            presentation_service, "build_presentation_model_manager", lambda: manager
        )
        patcher.start()
        self.addCleanup(patcher.stop)

    def use_retrieval(self, chunk_ids=(1, 2, 3), page: int = 7) -> None:
        async def fake_retrieve(*, search_query, **_kwargs):
            self.retrieval_queries.append(search_query)
            chunks = [
                {
                    "chunk_id": str(chunk_id),
                    "metadata": {
                        "doc_id": self.document.id,
                        "doc_name": self.document.name,
                        "page": page,
                    },
                }
                for chunk_id in chunk_ids
            ]
            texts = {str(chunk_id): f"текст фрагмента {chunk_id}" for chunk_id in chunk_ids}
            return chunks, texts

        patcher = patch.object(presentation_service, "retrieve_for_query", fake_retrieve)
        patcher.start()
        self.addCleanup(patcher.stop)

    # --- данные ---

    async def make_presentation(self, **overrides) -> Presentation:
        fields = {
            "notebook_id": self.notebook.id,
            "owner_id": self.user.id,
            "template_key": "classic",
            "language": "ru",
            "slide_count": DECK_SLIDES,
            "description": "Обзор льгот",
            "status": STATUS_QUEUED,
        }
        fields.update(overrides)
        return await self.seed(Presentation(**fields))

    async def run_one_job(self) -> Presentation:
        """Полный путь воркера: захват, генерация, запись исхода."""
        self.assertTrue(await self.worker._claim_and_process())
        row = await self.get_row(Presentation, self._presentation_id)
        self.assertIsNotNone(row)
        return row

    async def run_pipeline(self, *, timings=None):
        """Пайплайн напрямую, минуя цикл воркера.

        Нужен там, где проверяется САМ вызов, а не запись исхода: воркер
        заводит CallTimings внутри себя и наружу его не отдаёт, а отказ
        превращает в строку 'error'. Захват при этом настоящий — строка обязана
        быть в 'generating', иначе прогресс молча не пишется.
        """
        async with self.session_factory() as session:
            claimed = await presentation_service.PresentationsService.claim_next(
                session
            )
        self.assertEqual(claimed, self._presentation_id)
        return await presentation_service.generate_presentation(
            self._presentation_id, timings=timings
        )

    async def prepare(self, **overrides) -> None:
        self._presentation_id = (await self.make_presentation(**overrides)).id

    def storage_files(self) -> list[str]:
        return sorted(os.listdir(self.storage.name))


class SuccessfulGenerationTests(PresentationPipelineTestCase):
    async def asyncSetUp(self) -> None:
        await super().asyncSetUp()
        self.use_retrieval()
        self.use_model(
            [
                PLAN_JSON,
                slide_json("Кто имеет право", ["Первый факт", "Второй факт"], 1),
                slide_json("Как оформить", ["Третий факт", "Четвёртый факт"], 2),
                slide_json("Куда обращаться", ["Пятый факт", "Шестой факт"], 3),
            ]
        )
        await self.prepare()

    async def test_ready_row_points_to_a_real_file(self):
        row = await self.run_one_job()

        self.assertEqual(row.status, STATUS_READY)
        self.assertEqual(row.progress, 100)
        self.assertIsNone(row.error_code)
        self.assertTrue(os.path.exists(row.file_path))
        self.assertEqual(row.file_size, os.path.getsize(row.file_path))
        self.assertGreater(row.file_size, 0)

    async def test_no_temporary_file_survives(self):
        row = await self.run_one_job()
        self.assertEqual(self.storage_files(), [os.path.basename(row.file_path)])

    async def test_file_exists_before_the_row_promises_it(self):
        """Порядок «файл, потом commit» — зеркальный удалению.

        Обратный порядок оставил бы после падения между шагами status='ready'
        с путём в никуда: отказ на скачивании уже после «готово».
        """
        observed: list[bool] = []
        original = presentation_service.PresentationsService.mark_ready

        async def spy(session, presentation_id, *, file_path, file_size):
            observed.append(os.path.exists(file_path))
            return await original(
                session, presentation_id, file_path=file_path, file_size=file_size
            )

        with patch.object(
            presentation_service.PresentationsService, "mark_ready", spy
        ):
            await self.run_one_job()

        self.assertEqual(observed, [True])

    async def test_deck_has_the_ordered_slides_and_a_sources_slide(self):
        from pptx import Presentation as PptxPresentation

        row = await self.run_one_job()
        deck = PptxPresentation(row.file_path)
        texts = [
            "\n".join(
                shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
            )
            for slide in deck.slides
        ]

        self.assertEqual(len(texts), row.slide_count)
        self.assertIn("Налоговые льготы", texts[0])
        self.assertIn("Налоги", texts[0])  # имя блокнота
        self.assertIn("Кто имеет право", texts[1])
        self.assertIn("Как оформить", texts[2])
        self.assertIn("Куда обращаться", texts[3])
        self.assertIn("Источники", texts[-1])
        # Уникальные документы с именами и страницами.
        self.assertIn("Кодекс.pdf", texts[-1])
        self.assertIn("стр. 7", texts[-1])

    async def test_digest_of_written_bullets_reaches_the_next_slide(self):
        """Мера 2: слайд-вызов видит, что уже сказано на предыдущих."""
        await self.run_one_job()

        plan_call, first_slide, second_slide, _third_slide = self.model_calls
        self.assertNotIn("already_written", first_slide[1]["content"])
        self.assertIn("<already_written>", second_slide[1]["content"])
        self.assertIn("- Первый факт", second_slide[1]["content"])
        self.assertIn("- Второй факт", second_slide[1]["content"])
        # Заголовки предыдущих слайдов в дайджест не входят — только буллеты.
        self.assertNotIn("Кто имеет право\n</already_written>", second_slide[1]["content"])
        # И правило про два буллета стоит в системной части.
        self.assertIn("give exactly two bullets", second_slide[0]["content"])

    async def test_each_section_is_retrieved_by_its_own_query(self):
        await self.run_one_job()
        # Первый вызов — обзорный под план, дальше по одному на секцию.
        self.assertEqual(
            self.retrieval_queries,
            [
                "Обзор льгот",
                *[section["search_query"] for section in PLAN_SECTIONS],
            ],
        )

    async def test_journal_records_the_order_and_its_outcome(self):
        row = await self.run_one_job()

        logs = await self.all_rows(Log)
        self.assertEqual(len(logs), 1)
        entry = logs[0]
        self.assertEqual(entry.user_id, self.user.id)
        self.assertEqual(entry.notebook_id, self.notebook.id)
        self.assertIn("classic", entry.question)
        self.assertIn("ru", entry.question)
        self.assertIn(str(row.slide_count), entry.question)
        self.assertIn(STATUS_READY, entry.answer)
        self.assertIn("Кодекс.pdf", entry.sources or "")


class MissingSourcesTests(PresentationPipelineTestCase):
    async def test_no_indexed_sources_left(self):
        """Источники перечитываются в момент генерации, а не при заказе."""
        self.use_retrieval()
        self.use_model([])
        await self.prepare()

        async with self.session_factory() as session:
            document = await session.get(Document, self.document.id)
            document.status = "pending"
            session.add(document)
            await session.commit()

        row = await self.run_one_job()
        self.assertEqual(row.status, STATUS_ERROR)
        self.assertEqual(row.error_code, PresentationErrors.NO_SOURCES)
        self.assertTrue(row.error_text)
        # Ни одного файла: до рендера дело не дошло.
        self.assertEqual(self.storage_files(), [])

    async def test_documents_deleted_while_the_job_was_waiting(self):
        self.use_retrieval()
        self.use_model([])
        await self.prepare()

        async with self.session_factory() as session:
            document = await session.get(Document, self.document.id)
            await session.delete(document)
            await session.commit()

        row = await self.run_one_job()
        self.assertEqual(row.error_code, PresentationErrors.NO_SOURCES)

    async def test_journal_records_the_failure_too(self):
        self.use_retrieval()
        self.use_model([])
        await self.prepare()
        async with self.session_factory() as session:
            document = await session.get(Document, self.document.id)
            document.status = "pending"
            session.add(document)
            await session.commit()

        await self.run_one_job()

        logs = await self.all_rows(Log)
        self.assertEqual(len(logs), 1)
        self.assertIn(STATUS_ERROR, logs[0].answer)
        self.assertIn(PresentationErrors.NO_SOURCES, logs[0].answer)


class FailureTests(PresentationPipelineTestCase):
    async def test_unavailable_ollama_gets_its_own_code(self):
        self.use_retrieval()
        self.use_model(
            [ExternalServiceError("Ollama is unavailable", service="Ollama", status_code=503)]
        )
        await self.prepare()

        row = await self.run_one_job()
        self.assertEqual(row.status, STATUS_ERROR)
        self.assertEqual(row.error_code, PresentationErrors.OLLAMA_UNAVAILABLE)

    async def test_two_invalid_answers_fail_the_job(self):
        """Повтор один; после второго провала — честный отказ, а не починка."""
        self.use_retrieval()
        self.use_model(["не json", "тоже не json"])
        await self.prepare()

        row = await self.run_one_job()
        self.assertEqual(row.error_code, PresentationErrors.GENERATION_FAILED)
        self.assertEqual(len(self.model_calls), 2)
        # Вторая попытка получает исходный ответ и претензию валидатора.
        self.assertEqual(self.model_calls[1][-2]["content"], "не json")
        self.assertIn("rejected by the validator", self.model_calls[1][-1]["content"])

    async def test_retry_saves_a_recoverable_answer(self):
        self.use_retrieval()
        self.use_model(
            [
                "почти json",
                PLAN_JSON,
                slide_json("Кто имеет право", ["Факт", "Ещё факт"], 1),
                slide_json("Как оформить", ["Факт", "Ещё факт"], 2),
                slide_json("Куда обращаться", ["Факт", "Ещё факт"], 3),
            ]
        )
        await self.prepare()

        row = await self.run_one_job()
        self.assertEqual(row.status, STATUS_READY)

    async def test_an_invented_citation_is_refused_and_the_retry_saves_the_deck(self):
        """Ссылка на чанк, которого модели не показывали, — не слайд.

        Проверка «цитата лежит внутри выданного набора» есть в схеме
        (tests/test_presentation_schemas.py), но там она получает
        allowed_citations из фикстуры. Здесь набор собирает сам пайплайн — из
        того, что вернул ретривал, — и проверяется вся дорога: выдуманный
        chunk_id отвергается, претензия валидатора с перечнем разрешённых
        значений уезжает в повторный промпт, и исправленный ответ доводит
        колоду до конца.
        """
        self.use_retrieval(chunk_ids=(1, 2, 3))
        self.use_model(
            [
                PLAN_JSON,
                # Чанка 99 в контексте не было: ретривал отдал 1, 2 и 3.
                slide_json("Кто имеет право", ["Факт", "Ещё факт"], 99),
                slide_json("Кто имеет право", ["Факт", "Ещё факт"], 1),
                slide_json("Как оформить", ["Факт", "Ещё факт"], 2),
                slide_json("Куда обращаться", ["Факт", "Ещё факт"], 3),
            ]
        )
        await self.prepare()

        row = await self.run_one_job()

        self.assertEqual(row.status, STATUS_READY)
        # Повтор получил и отвергнутый ответ, и предметную претензию: без
        # перечня разрешённых chunk_id модель повторила бы ту же выдумку.
        retry_prompt = self.model_calls[2][-1]["content"]
        self.assertIn("rejected by the validator", retry_prompt)
        self.assertIn("'99'", retry_prompt)

    async def test_two_invented_citations_fail_the_job_without_a_file(self):
        """Выдуманный источник не «чинится» выбрасыванием цитаты.

        Слайд недействителен целиком: неизвестно, какое из его утверждений
        опиралось на несуществующий фрагмент. Поэтому исход — отказ, а не
        колода с правдоподобным списком источников; это же и есть то, что
        обесценивает инъекцию через описание заказа.
        """
        self.use_retrieval(chunk_ids=(1, 2, 3))
        self.use_model(
            [
                PLAN_JSON,
                slide_json("Кто имеет право", ["Факт", "Ещё факт"], 77),
                slide_json("Кто имеет право", ["Факт", "Ещё факт"], 88),
            ]
        )
        await self.prepare()

        row = await self.run_one_job()

        self.assertEqual(row.status, STATUS_ERROR)
        self.assertEqual(row.error_code, PresentationErrors.GENERATION_FAILED)
        self.assertEqual(self.storage_files(), [])
        # Ровно две попытки на слайд: план плюс два слайд-вызова.
        self.assertEqual(len(self.model_calls), 3)

    async def test_a_hanging_model_times_out_and_leaves_nothing_behind(self):
        """Ollama приняла запрос и не ответила: вызов снимает СВОЙ таймаут.

        Главное здесь — ЧЕЙ таймаут сработал. Потолок джобы не тронут и
        остаётся настоящим (для колоды из SLIDE_COUNT_MIN слайдов это
        presentation_job_timeout(5) = 2700 с): если бы зависший вызов ждал
        его, тест не уложился бы и в сорок минут. Снят же заказ за доли
        секунды — то есть по бюджету ОДНОГО вызова. Ровно это и было целью
        переноса: повисший вызов на втором слайде из пятнадцати не должен
        убивать заказ через двадцать минут.

        Заодно проверяется то, чего не проверяет тест на подменённом пайплайне
        (tests/test_presentation_queue_db.py): здесь пайплайн настоящий, и
        снимают его посреди вызова модели — внутри открытой сессии, после
        ретривала и до всякого файла. После снятого заказа на диске не осталось
        НИЧЕГО — ни колоды, ни `.tmp-*`, — а строка не обещает файла.

        Бюджет: подготовка до первого обращения к модели — несколько запросов к
        PostgreSQL, единицы миллисекунд; таймаут вызова взят на два порядка
        больше, поэтому «не успели дойти до модели» здесь не случайность, а
        поломка, и её ловит проверка entered ниже.
        """
        self.use_retrieval()
        entered: list[str] = []

        class HangingModel:
            async def chat(self, *, model=None, messages=None, num_ctx=None) -> str:
                entered.append("chat")
                await asyncio.sleep(30)
                raise AssertionError("висящий вызов не должен завершаться")

        patcher = patch.object(
            presentation_service,
            "build_presentation_model_manager",
            lambda: HangingModel(),
        )
        patcher.start()
        self.addCleanup(patcher.stop)
        await self.prepare()

        started = perf_counter()
        # Подменённая модель не ходит по HTTP, значит клиентскому таймауту
        # (LLM_CALL_TIMEOUT) сработать не на чем — здесь проверяется именно
        # СТРАХОВКА второго эшелона, wait_for. В бою у неё та же роль:
        # зависание, которого клиент не видит по построению.
        with patch.object(presentation_service, "LLM_CALL_WATCHDOG_TIMEOUT", 1.0):
            row = await self.run_one_job()
        elapsed = perf_counter() - started

        self.assertEqual(entered, ["chat"], "до вызова модели дело не дошло")
        # Порог с большим запасом к бюджету вызова (1 с) и с большим запасом
        # ВНИЗ к тому, сколько ждал бы заказ, если бы его снимал не он:
        # висящий вызов заготовлен на 30 с, потолок джобы — 2700 с.
        self.assertLess(
            elapsed,
            10.0,
            "заказ ждал общего потолка джобы "
            f"({presentation_job_timeout(DECK_SLIDES):.0f} с), а не потолка вызова",
        )
        self.assertEqual(row.status, STATUS_ERROR)
        self.assertEqual(row.error_code, PresentationErrors.GENERATION_TIMEOUT)
        # error_text называет СТАДИЮ: «не уложились во время» без стадии
        # одинаково описывает и первую минуту, и предпоследнюю секцию.
        self.assertIn(presentation_service.STAGE_PLAN, row.error_text)
        # Причина понятна пользователю и не тащит с собой ни трейсбека, ни
        # путей на сервере (error_text уходит клиенту вместе со статусом).
        self.assertNotIn("Traceback", row.error_text)
        self.assertNotIn(self.storage.name, row.error_text)
        # И главное: диск чист, а строка ничего не обещает.
        self.assertIsNone(row.file_path)
        self.assertEqual(self.storage_files(), [])

    async def test_a_hanging_slide_call_names_the_slide_in_the_error(self):
        """Отказ по таймауту называет не только стадию, но и номер слайда.

        План прошёл, повис второй слайд-вызов. Пользователь видит строку с
        причиной рядом со статусом, и «слайд 2 из 3» — единственное, что
        отличает её от такой же строки про первую минуту генерации.
        """
        self.use_retrieval()
        answers = [PLAN_JSON, slide_json("Кто имеет право", ["Факт", "Ещё"], 1)]

        class SlowOnTheSecondSlide:
            async def chat(self, *, model=None, messages=None, num_ctx=None) -> str:
                if answers:
                    return answers.pop(0)
                await asyncio.sleep(30)
                raise AssertionError("висящий вызов не должен завершаться")

        patcher = patch.object(
            presentation_service,
            "build_presentation_model_manager",
            lambda: SlowOnTheSecondSlide(),
        )
        patcher.start()
        self.addCleanup(patcher.stop)
        await self.prepare()

        with patch.object(presentation_service, "LLM_CALL_WATCHDOG_TIMEOUT", 1.0):
            row = await self.run_one_job()

        self.assertEqual(row.error_code, PresentationErrors.GENERATION_TIMEOUT)
        self.assertIn("слайд 2 из 3", row.error_text)
        self.assertIn(PLAN_SECTIONS[1]["heading"], row.error_text)
        self.assertEqual(self.storage_files(), [])

    async def test_the_retry_gets_a_budget_of_its_own(self):
        """Повтор — отдельный вызов, а не остаток бюджета первой попытки.

        Он получает исходный промпт, отвергнутый ответ и претензию валидатора,
        то есть генерирует ВЕСЬ ответ заново и стоит примерно столько же.
        Общий на две попытки бюджет означал бы, что медленная первая попытка
        съедает время второй: заказ падал бы по таймауту на попытке, которая
        сама по себе была здоровой.

        Бюджет теста: потолок вызова 0.5 с, первая попытка занимает 0.4 с из
        них. При ОБЩЕМ бюджете второй попытке осталось бы 0.1 с; проверка ниже
        требует от неё хотя бы 0.3 с — то есть промах втрое, а не на границе
        точности таймера.
        """
        self.use_retrieval()
        call_budget = 0.5
        attempts: list[int] = []

        class SlowThenHanging:
            async def chat(self, *, model=None, messages=None, num_ctx=None) -> str:
                attempts.append(len(attempts) + 1)
                if len(attempts) == 1:
                    await asyncio.sleep(call_budget * 0.8)
                    return "это не JSON"  # валидатор отвергнет -> повтор
                await asyncio.sleep(30)
                raise AssertionError("висящий вызов не должен завершаться")

        patcher = patch.object(
            presentation_service,
            "build_presentation_model_manager",
            lambda: SlowThenHanging(),
        )
        patcher.start()
        self.addCleanup(patcher.stop)
        await self.prepare()

        timings = presentation_service.CallTimings()
        with patch.object(
            presentation_service, "LLM_CALL_WATCHDOG_TIMEOUT", call_budget
        ):
            with self.assertRaises(
                presentation_service.PresentationGenerationError
            ) as caught:
                await self.run_pipeline(timings=timings)

        self.assertEqual(attempts, [1, 2], "повторной попытки не было")
        self.assertEqual(
            caught.exception.error_code, PresentationErrors.GENERATION_TIMEOUT
        )
        self.assertEqual(timings.retries, 1)
        self.assertGreaterEqual(
            timings.durations[1],
            call_budget * 0.6,
            "повтор получил остаток бюджета первой попытки, а не свой",
        )

    async def test_call_timings_are_collected_for_the_statistics_line(self):
        """Каждый вызов модели попадает в CallTimings, и план отличим от слайда."""
        self.use_retrieval()
        self.use_model(
            [
                PLAN_JSON,
                slide_json("Кто имеет право", ["Факт", "Ещё факт"], 1),
                slide_json("Как оформить", ["Факт", "Ещё факт"], 2),
                slide_json("Куда обращаться", ["Факт", "Ещё факт"], 3),
            ]
        )
        await self.prepare()

        timings = presentation_service.CallTimings()
        await self.run_pipeline(timings=timings)

        self.assertEqual(timings.plan_calls, 1)
        self.assertEqual(timings.slide_calls, len(PLAN_SECTIONS))
        self.assertEqual(timings.retries, 0)
        self.assertEqual(len(timings.durations), 1 + len(PLAN_SECTIONS))
        summary = timings.summary()
        self.assertIn("p50", summary)
        self.assertIn("p90", summary)

    async def test_broken_render_leaves_no_files_behind(self):
        self.use_retrieval()
        self.use_model(
            [
                PLAN_JSON,
                slide_json("Кто имеет право", ["Факт", "Ещё факт"], 1),
                slide_json("Как оформить", ["Факт", "Ещё факт"], 2),
                slide_json("Куда обращаться", ["Факт", "Ещё факт"], 3),
            ]
        )
        await self.prepare()

        with patch.object(
            presentation_service, "render_presentation", side_effect=OSError("диск полон")
        ):
            row = await self.run_one_job()

        self.assertEqual(row.status, STATUS_ERROR)
        self.assertEqual(row.error_code, PresentationErrors.GENERATION_FAILED)
        self.assertIsNone(row.file_path)
        self.assertEqual(self.storage_files(), [])

    async def test_row_deleted_while_generating_is_not_resurrected(self):
        """Отменённый заказ не должен вернуться в базу строкой 'ready'."""
        self.use_retrieval()
        self.use_model([])
        await self.prepare()

        async with self.session_factory() as session:
            row = await session.get(Presentation, self._presentation_id)
            row.status = STATUS_GENERATING
            session.add(row)
            await session.commit()
            await session.delete(row)
            await session.commit()

        # Захватывать нечего: строки нет.
        self.assertFalse(await self.worker._claim_and_process())
        self.assertFalse(await self.exists(Presentation, self._presentation_id))


if __name__ == "__main__":
    unittest.main()
