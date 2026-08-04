"""Рендер колоды: устойчивость к длинным строкам и состав файла.

Ни базы, ни модели: рендерер получает готовую структуру и отдаёт файл. Здесь
проверяется главное его правило — длинная строка ОБРЕЗАЕТСЯ, а не роняет
рендер. Уронить генерацию на имени документа означало бы выбросить всю уже
сделанную работу модели, которая стоила пользователю минуты ожидания.

Реестр шаблонов здесь намеренно не нужен: рендер обязан работать и с ним, и
без него (до появления templates.py он был единственным способом отдать
пользователю файл), поэтому основной прогон идёт по несуществующему ключу.
"""

import os
import sys
import tempfile
import unittest
from datetime import datetime

sys.path.insert(0, os.path.dirname(os.path.abspath(__file__)))

from app.modules.presentations.llm_schemas import PresentationSlide  # noqa: E402
from app.modules.presentations.renderer import (  # noqa: E402
    MAX_SOURCE_LINES,
    SOURCE_LINE_MAX_CHARS,
    RenderedSource,
    build_source_lines,
    fit,
    format_source_line,
)
from app.modules.presentations.renderer import render_presentation  # noqa: E402


class FitTests(unittest.TestCase):
    def test_short_string_is_untouched(self):
        self.assertEqual(fit("Ставка НДС", 80), "Ставка НДС")

    def test_long_string_is_cut_with_an_ellipsis(self):
        result = fit("я" * 500, 80)
        self.assertEqual(len(result), 80)
        self.assertTrue(result.endswith("…"))

    def test_whitespace_is_collapsed(self):
        # Перевод строки внутри имени документа ломал бы список источников.
        self.assertEqual(fit(" Налоговый\n  кодекс ", 80), "Налоговый кодекс")

    def test_empty_value_is_allowed(self):
        self.assertEqual(fit(None, 80), "")


class SourceLineTests(unittest.TestCase):
    def test_pages_are_unique_and_sorted(self):
        line = format_source_line(
            RenderedSource(source_id=1, name="Кодекс.pdf", pages=[5, 1, 5, 3]), "ru"
        )
        self.assertEqual(line, "Кодекс.pdf — стр. 1, 3, 5")

    def test_document_without_pages_keeps_only_its_name(self):
        line = format_source_line(
            RenderedSource(source_id=1, name="Кодекс.pdf", pages=[]), "ru"
        )
        self.assertEqual(line, "Кодекс.pdf")

    def test_tajik_label(self):
        line = format_source_line(
            RenderedSource(source_id=1, name="Кодекс.pdf", pages=[2]), "tj"
        )
        self.assertEqual(line, "Кодекс.pdf — саҳ. 2")

    def test_nameless_document_falls_back_to_its_id(self):
        line = format_source_line(RenderedSource(source_id=42, name="", pages=[]), "ru")
        self.assertEqual(line, "#42")

    def test_long_name_is_cut_and_pages_survive(self):
        line = format_source_line(
            RenderedSource(source_id=1, name="и" * 400, pages=[7]), "ru"
        )
        self.assertLessEqual(len(line), SOURCE_LINE_MAX_CHARS)
        self.assertTrue(line.endswith("стр. 7"))

    def test_overflowing_list_says_how_many_are_hidden(self):
        sources = [
            RenderedSource(source_id=index, name=f"Документ {index}.pdf", pages=[1])
            for index in range(MAX_SOURCE_LINES + 5)
        ]
        lines = build_source_lines(sources, "ru")
        self.assertEqual(len(lines), MAX_SOURCE_LINES + 1)
        self.assertEqual(lines[-1], "… +5")

    def test_empty_list_still_gives_a_line(self):
        self.assertEqual(build_source_lines([], "ru"), ["—"])


class RenderTests(unittest.TestCase):
    def setUp(self):
        self._dir = tempfile.TemporaryDirectory()
        self.addCleanup(self._dir.cleanup)
        self.path = os.path.join(self._dir.name, "deck.pptx")

    def render(self, **overrides):
        params = {
            "title": "Налоговые льготы",
            "slides": [
                PresentationSlide(
                    heading="Кто имеет право",
                    bullets=["Первый факт", "Второй факт"],
                    citations=[{"source_id": 1, "chunk_id": 10}],
                )
            ],
            "sources": [
                RenderedSource(source_id=1, name="Кодекс.pdf", pages=[3])
            ],
            "language": "ru",
            "template_key": "no-such-template",
            "notebook_name": "Налоги",
            "created_at": datetime(2026, 8, 4, 9, 15),
            "output_path": self.path,
        }
        params.update(overrides)
        render_presentation(**params)
        from pptx import Presentation

        deck = Presentation(self.path)
        return [
            "\n".join(
                shape.text_frame.text for shape in slide.shapes if shape.has_text_frame
            )
            for slide in deck.slides
        ]

    def test_deck_shape_is_title_content_sources(self):
        texts = self.render()
        self.assertEqual(len(texts), 3)
        self.assertIn("Налоговые льготы", texts[0])
        self.assertIn("Налоги", texts[0])
        self.assertIn("04.08.2026", texts[0])
        self.assertIn("Кто имеет право", texts[1])
        self.assertIn("Первый факт", texts[1])
        self.assertIn("Источники", texts[2])
        self.assertIn("Кодекс.pdf — стр. 3", texts[2])

    def test_absurdly_long_values_do_not_break_the_render(self):
        texts = self.render(
            title="Название " * 200,
            notebook_name="Блокнот " * 200,
            sources=[RenderedSource(source_id=1, name="и" * 900, pages=[1])],
        )
        self.assertEqual(len(texts), 3)
        self.assertTrue(os.path.getsize(self.path) > 0)
        for text in texts:
            for line in text.splitlines():
                self.assertLessEqual(len(line), 500)

    def test_special_characters_do_not_break_the_render(self):
        """Спецсимволы — вторая половина правила устойчивости.

        Внутри pptx лежит XML, и текст туда попадает как есть. Опасны три
        разных вида символов, и приходят они с разных сторон:

          * `&`, `<`, `>`, кавычки — из имени документа и из ответа модели.
            Неэкранированные, они дали бы битый XML, то есть файл, который
            PowerPoint отказывается открыть, — худший исход из возможных:
            строка 'ready', размер больше нуля, а колоды нет;
          * управляющие символы (`\\x00`, `\\x07`) — из ответа модели: JSON
            умеет их записывать (`\\u0007`), а XML 1.0 не допускает вовсе;
          * эмодзи и письменности вне BMP — из документов пользователя.

        Проверяется не «как именно», а результат: файл собран, открывается и
        текст на месте. Экранированием занимается python-pptx, и подменять его
        своей чисткой было бы вторым правилом на ту же тему.
        """
        nasty = 'R&D <tag> "кавычки" & 🧾 متن'
        texts = self.render(
            title=nasty,
            slides=[
                PresentationSlide(
                    heading=nasty,
                    bullets=["Строка\x07с\x00управляющими", "Второй факт"],
                    citations=[{"source_id": 1, "chunk_id": 10}],
                )
            ],
            sources=[RenderedSource(source_id=1, name=nasty, pages=[3])],
        )

        self.assertEqual(len(texts), 3)
        self.assertTrue(os.path.getsize(self.path) > 0)
        # Файл читается обратно (значит, XML целый), и текст в нём тот самый.
        self.assertIn("R&D <tag>", texts[0])
        self.assertIn("🧾", texts[1])
        self.assertIn(nasty, texts[2])
        # Управляющие символы доехали, не уронив рендер: python-pptx
        # записывает их как _xNNNN_, и это его дело, а не наше.
        self.assertIn("Второй факт", texts[1])

    def test_tajik_deck_uses_tajik_captions(self):
        texts = self.render(language="tj")
        self.assertIn("Манбаъҳо", texts[2])
        self.assertIn("саҳ.", texts[2])


if __name__ == "__main__":
    unittest.main()
