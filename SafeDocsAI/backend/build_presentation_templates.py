"""Сборка файлов шаблонов презентаций для templates/presentations/.

Зачем скрипт, а не «положили pptx руками». Шаблоны v1 сделаны кодом, потому
что дизайнерской работы на этапе 1 нет, а рендереру нужны рабочие файлы уже
сейчас. Скрипт фиксирует, ЧЕМ именно отличаются три комплекта (палитра,
гарнитура, композиция) и позволяет пересобрать их после правки — иначе первый
же вопрос «почему в тёмном шаблоне заголовок серый» упирался бы в двоичный
файл, который нечем прочитать.

Когда придёт настоящий дизайн, скрипт заменяется присланными pptx: реестр
(app/modules/presentations/templates.py) читает манифест и файлы, а не этот
модуль, и о его существовании не знает.

Запуск:

    ./venv/bin/python build_presentation_templates.py

Перезаписывает templates/presentations/*.pptx и *.png; manifest.json НЕ трогает
— он правится руками, потому что там же живут ключи и названия на двух языках.
"""

from __future__ import annotations

import sys
from dataclasses import dataclass
from pathlib import Path

from PIL import Image, ImageDraw, ImageFont
from pptx import Presentation
from pptx.oxml.ns import nsdecls, qn
from pptx.util import Emu, Inches
from lxml import etree

# Шаблоны 16:9: 4:3 из умолчаний python-pptx на проекторе даёт поля по бокам,
# а презентация здесь — экспортный артефакт, который смотрят с ноутбука.
SLIDE_WIDTH = Inches(13.333)
SLIDE_HEIGHT = Inches(7.5)

# Роли layout'ов ровно те, что перечислены в манифесте и в реестре. Порядок
# важен: после удаления лишних layout'ов индекс роли равен её позиции здесь.
LAYOUT_ROLES = ("title", "bullets", "section", "sources")

# Какие layout'ы стандартного шаблона python-pptx под какую роль берём.
# 0 Title Slide, 1 Title and Content, 2 Section Header, 3 Two Content.
# «Два столбца» переделывается в «Источники»: второй блок текста удаляется,
# первый растягивается на всю ширину. Отдельный layout под источники нужен,
# чтобы дизайнер потом оформил их иначе, чем обычные буллиты, не трогая код.
SOURCE_LAYOUT_INDEXES = (0, 1, 2, 3)

PREVIEW_SIZE = (1280, 720)
_FONT_DIR = Path("/usr/share/fonts/truetype/dejavu")


@dataclass(frozen=True)
class Palette:
    """Оформление одного шаблона.

    Все цвета — hex без решётки: в таком виде их принимает и srgbClr в pptx, и
    Pillow при отрисовке превью, поэтому палитра для файла и для картинки одна
    и та же и разъехаться не может.
    """

    key: str
    title_ru: str
    background: str
    surface: str  # подложка блока текста на превью
    accent: str
    heading_color: str
    body_color: str
    muted_color: str
    heading_font: str
    body_font: str
    # Полоса-акцент у титульного слайда: False — акцент даёт тонкая линия.
    accent_bar: bool


PALETTES = (
    Palette(
        key="classic",
        title_ru="Классический",
        background="FFFFFF",
        surface="F2F5FA",
        accent="1F3A63",
        heading_color="1F3A63",
        body_color="2E3440",
        muted_color="6B7684",
        heading_font="Calibri",
        body_font="Calibri",
        accent_bar=True,
    ),
    Palette(
        key="contrast",
        title_ru="Контрастный",
        background="101820",
        surface="1B2530",
        accent="F2A900",
        heading_color="FFFFFF",
        body_color="D8DEE9",
        muted_color="8B97A6",
        heading_font="Arial",
        body_font="Arial",
        accent_bar=True,
    ),
    Palette(
        key="minimal",
        title_ru="Минималистичный",
        background="FFFFFF",
        surface="FFFFFF",
        accent="9AA0A6",
        heading_color="111111",
        body_color="333333",
        muted_color="8A8A8A",
        heading_font="Georgia",
        body_font="Georgia",
        accent_bar=False,
    ),
)


# --- Правка XML layout'ов ------------------------------------------------
#
# python-pptx умеет двигать плейсхолдеры, но не умеет задавать их ТИПОГРАФИКУ
# так, чтобы она досталась слайду: шрифт, выставленный на тексте-подсказке
# layout'а, на слайд не наследуется. Наследуется список стилей уровней
# (a:lstStyle) внутри плейсхолдера, поэтому пишем именно его.


def _set_list_style(placeholder, levels: list[str]) -> None:
    """Заменить a:lstStyle плейсхолдера на переданные уровни (lvl1..lvlN)."""
    txBody = placeholder.text_frame._txBody
    for existing in txBody.findall(qn("a:lstStyle")):
        txBody.remove(existing)
    xml = f'<a:lstStyle {nsdecls("a")}>{"".join(levels)}</a:lstStyle>'
    lstStyle = etree.fromstring(xml)
    # Порядок элементов в a:txBody фиксирован схемой: bodyPr, lstStyle, p+.
    bodyPr = txBody.find(qn("a:bodyPr"))
    bodyPr.addnext(lstStyle)


def _heading_level(palette: Palette, size_pt: int, *, color: str | None = None) -> str:
    return (
        f'<a:lvl1pPr marL="0" indent="0"><a:buNone/>'
        f'<a:defRPr sz="{size_pt * 100}" b="1">'
        f'<a:solidFill><a:srgbClr val="{color or palette.heading_color}"/></a:solidFill>'
        f'<a:latin typeface="{palette.heading_font}"/>'
        f"</a:defRPr></a:lvl1pPr>"
    )


def _plain_level(palette: Palette, size_pt: int, *, color: str | None = None) -> str:
    return (
        f'<a:lvl1pPr marL="0" indent="0"><a:buNone/>'
        f'<a:defRPr sz="{size_pt * 100}">'
        f'<a:solidFill><a:srgbClr val="{color or palette.muted_color}"/></a:solidFill>'
        f'<a:latin typeface="{palette.body_font}"/>'
        f"</a:defRPr></a:lvl1pPr>"
    )


def _bullet_levels(palette: Palette) -> list[str]:
    """Два уровня буллитов: слайд-контент и вложенный пункт."""
    levels = []
    for index, (size_pt, indent_in, char) in enumerate(
        ((20, 0.0, "•"), (16, 0.35, "–")), start=1
    ):
        mar = int(Inches(indent_in + 0.3))
        levels.append(
            f'<a:lvl{index}pPr marL="{mar}" indent="-{int(Inches(0.3))}">'
            f'<a:spcBef><a:spcPts val="600"/></a:spcBef>'
            f'<a:buClr><a:srgbClr val="{palette.accent}"/></a:buClr>'
            f'<a:buChar char="{char}"/>'
            f'<a:defRPr sz="{size_pt * 100}">'
            f'<a:solidFill><a:srgbClr val="{palette.body_color}"/></a:solidFill>'
            f'<a:latin typeface="{palette.body_font}"/>'
            f"</a:defRPr></a:lvl{index}pPr>"
        )
    return levels


def _drop_placeholder(layout, idx: int) -> None:
    """Убрать плейсхолдер по idx, если он есть."""
    for placeholder in list(layout.placeholders):
        if placeholder.placeholder_format.idx == idx:
            element = placeholder.element
            element.getparent().remove(element)
            return


def _place(placeholder, left, top, width, height) -> None:
    placeholder.left = Emu(int(left))
    placeholder.top = Emu(int(top))
    placeholder.width = Emu(int(width))
    placeholder.height = Emu(int(height))


def _set_background(container, color: str) -> None:
    fill = container.background.fill
    fill.solid()
    fill.fore_color.rgb = _rgb(color)


def _rgb(color: str):
    from pptx.dml.color import RGBColor

    return RGBColor.from_string(color)


def _add_accent_bar(layout, palette: Palette, left, top, width, height) -> None:
    """Прямоугольник-акцент. У minimal его нет — там роль акцента у линии.

    add_shape есть у слайда, но не у layout'а (LayoutShapes умеет только
    читать), поэтому фигура собирается тем же элементом XML, которым её
    собирает сам python-pptx, и кладётся в дерево фигур layout'а.
    """
    from pptx.oxml.shapes.autoshape import CT_Shape
    from pptx.shapes.autoshape import Shape

    spTree = layout.shapes._spTree
    shape_id = max(
        (int(el.get("id")) for el in spTree.iter() if el.tag.endswith("}cNvPr")),
        default=1,
    ) + 1
    sp = CT_Shape.new_autoshape_sp(
        shape_id, "Accent Bar", "rect", int(left), int(top), int(width), int(height)
    )
    spTree.append(sp)
    shape = Shape(sp, layout.shapes)
    shape.fill.solid()
    shape.fill.fore_color.rgb = _rgb(palette.accent)
    shape.line.fill.background()
    shape.shadow.inherit = False


def _build_pptx(palette: Palette, target: Path) -> None:
    presentation = Presentation()
    presentation.slide_width = SLIDE_WIDTH
    presentation.slide_height = SLIDE_HEIGHT

    master = presentation.slide_masters[0]
    _set_background(master, palette.background)

    # Лишние layout'ы удаляются, чтобы индексы в манифесте оставались
    # осмысленными: 0/1/2/3 — ровно четыре роли, а не «где-то среди
    # одиннадцати заготовок Microsoft».
    keep = [presentation.slide_layouts[i] for i in SOURCE_LAYOUT_INDEXES]
    for layout in list(presentation.slide_layouts):
        if layout not in keep:
            master.slide_layouts.remove(layout)

    title_layout, bullets_layout, section_layout, sources_layout = keep
    title_layout.name = "Title"
    bullets_layout.name = "Bullets"
    section_layout.name = "Section"
    sources_layout.name = "Sources"

    margin = Inches(0.9)
    content_width = SLIDE_WIDTH - 2 * margin

    for layout in keep:
        _set_background(layout, palette.background)
        # Дата и подпись не нужны: ни то, ни другое рендерер не заполняет, а
        # пустые плейсхолдеры видны в редакторе и мешают правке файла руками.
        _drop_placeholder(layout, 10)
        _drop_placeholder(layout, 11)

    # --- Титул -----------------------------------------------------------
    if palette.accent_bar:
        _add_accent_bar(
            title_layout, palette, margin, Inches(2.3), Inches(1.6), Inches(0.11)
        )
    title_ph = title_layout.placeholders[0]
    _place(title_ph, margin, Inches(2.7), content_width, Inches(2.0))
    _set_list_style(title_ph, [_heading_level(palette, 40)])
    subtitle_ph = title_layout.placeholders[1]
    _place(subtitle_ph, margin, Inches(4.8), content_width, Inches(0.9))
    _set_list_style(subtitle_ph, [_plain_level(palette, 18)])
    _drop_placeholder(title_layout, 12)

    # --- Раздел ----------------------------------------------------------
    if palette.accent_bar:
        _add_accent_bar(
            section_layout, palette, margin, Inches(2.9), Inches(0.9), Inches(0.11)
        )
    section_title = section_layout.placeholders[0]
    _place(section_title, margin, Inches(3.2), content_width, Inches(1.4))
    _set_list_style(section_title, [_heading_level(palette, 32)])
    # Подзаголовок раздела остаётся: рендерер может положить туда номер части,
    # а пустым он на слайде не виден.
    section_text = section_layout.placeholders[1]
    _place(section_text, margin, Inches(4.5), content_width, Inches(0.8))
    _set_list_style(section_text, [_plain_level(palette, 16)])
    _drop_placeholder(section_layout, 12)

    # --- Буллиты ---------------------------------------------------------
    bullets_title = bullets_layout.placeholders[0]
    _place(bullets_title, margin, Inches(0.7), content_width, Inches(1.0))
    _set_list_style(bullets_title, [_heading_level(palette, 28)])
    bullets_body = bullets_layout.placeholders[1]
    _place(bullets_body, margin, Inches(1.9), content_width, Inches(4.7))
    _set_list_style(bullets_body, _bullet_levels(palette))
    _slide_number(bullets_layout, palette)

    # --- Источники -------------------------------------------------------
    sources_title = sources_layout.placeholders[0]
    _place(sources_title, margin, Inches(0.7), content_width, Inches(1.0))
    _set_list_style(sources_title, [_heading_level(palette, 28)])
    # Второй столбец «Two Content» здесь не нужен: список источников — один
    # поток строк, и разрывать его пополам нечем.
    _drop_placeholder(sources_layout, 2)
    sources_body = sources_layout.placeholders[1]
    _place(sources_body, margin, Inches(1.9), content_width, Inches(4.7))
    _set_list_style(
        sources_body,
        [
            f'<a:lvl1pPr marL="0" indent="0"><a:buNone/>'
            f'<a:spcBef><a:spcPts val="500"/></a:spcBef>'
            f'<a:defRPr sz="1400">'
            f'<a:solidFill><a:srgbClr val="{palette.muted_color}"/></a:solidFill>'
            f'<a:latin typeface="{palette.body_font}"/>'
            f"</a:defRPr></a:lvl1pPr>"
        ],
    )
    _slide_number(sources_layout, palette)

    target.parent.mkdir(parents=True, exist_ok=True)
    presentation.save(str(target))


def _slide_number(layout, palette: Palette) -> None:
    """Номер слайда внизу справа; на титуле и разделителе он не нужен."""
    for placeholder in layout.placeholders:
        if placeholder.placeholder_format.idx == 12:
            _place(
                placeholder,
                SLIDE_WIDTH - Inches(1.6),
                SLIDE_HEIGHT - Inches(0.75),
                Inches(0.9),
                Inches(0.4),
            )
            _set_list_style(placeholder, [_plain_level(palette, 11)])
            return


# --- Превью --------------------------------------------------------------
#
# Настоящего рендера pptx в картинку в образе нет (LibreOffice не ставится
# ради одного превью), поэтому картинка рисуется из ТОЙ ЖЕ палитры и по той же
# композиции, что и титульный слайд: это схема оформления, а не скриншот.
# Отличать их честнее, чем выдавать одно за другое.


def _font(size: int, *, bold: bool = False) -> ImageFont.FreeTypeFont:
    name = "DejaVuSans-Bold.ttf" if bold else "DejaVuSans.ttf"
    path = _FONT_DIR / name
    if path.exists():
        return ImageFont.truetype(str(path), size)
    return ImageFont.load_default(size)


def _build_preview(palette: Palette, target: Path) -> None:
    width, height = PREVIEW_SIZE
    image = Image.new("RGB", PREVIEW_SIZE, f"#{palette.background}")
    draw = ImageDraw.Draw(image)

    margin = 86
    if palette.accent_bar:
        draw.rectangle(
            [margin, 250, margin + 154, 260], fill=f"#{palette.accent}"
        )
    else:
        draw.line([margin, 255, width - margin, 255], fill=f"#{palette.accent}", width=2)

    draw.text(
        (margin, 292),
        palette.title_ru,
        font=_font(58, bold=True),
        fill=f"#{palette.heading_color}",
    )
    draw.text(
        (margin, 386),
        "SafeDocsAI · шаблон презентации",
        font=_font(24),
        fill=f"#{palette.muted_color}",
    )

    # Полоска-«контент» снизу: показывает подложку и цвет буллитов.
    block_top = 486
    draw.rectangle(
        [margin, block_top, width - margin, height - 70],
        fill=f"#{palette.surface}",
    )
    for row in range(3):
        y = block_top + 34 + row * 40
        draw.rectangle([margin + 28, y, margin + 40, y + 12], fill=f"#{palette.accent}")
        draw.rectangle(
            [margin + 60, y + 2, width - margin - 120 - row * 90, y + 10],
            fill=f"#{palette.body_color}",
        )

    target.parent.mkdir(parents=True, exist_ok=True)
    image.save(str(target), format="PNG")


def main() -> int:
    backend_dir = Path(__file__).resolve().parent
    directory = backend_dir / "templates" / "presentations"
    for palette in PALETTES:
        pptx_path = directory / f"{palette.key}.pptx"
        preview_path = directory / f"{palette.key}.png"
        _build_pptx(palette, pptx_path)
        _build_preview(palette, preview_path)
        print(f"{palette.key}: {pptx_path.name}, {preview_path.name}")
    return 0


if __name__ == "__main__":
    sys.exit(main())
