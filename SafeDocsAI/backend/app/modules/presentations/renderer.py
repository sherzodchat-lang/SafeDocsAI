"""Сборка .pptx из провалидированной структуры.

Рендер — единственное место пайплайна, где работает python-pptx, и он
СИНХРОННЫЙ: библиотека распаковывает и пакует zip, а это блокирует поток.
Поэтому звать его можно только через run_in_threadpool (см. service.py) —
тот же дефект в проекте уже чинили для ChromaDB.

Модуль ничего не знает ни о базе, ни о модели: на вход — план, слайды и
источники, на выход — файл по указанному пути. Из-за этого его можно проверить
без PostgreSQL и без Ollama.

Правило устойчивости: длинная строка обрезается, а не роняет рендер. Схема
ответа модели уже держит границы (heading 80, буллет 200), но в рендер
приходят ещё и имена документов, имя блокнота и заголовок плана — их длину
никто не обещал.
"""

import logging
from dataclasses import dataclass, field
from datetime import datetime
from typing import Any

from app.modules.presentations.constants import (
    PAGES_LABEL,
    SOURCES_HEADING,
    SLIDES_WORD,
)
from app.modules.presentations.llm_schemas import (
    PresentationSlide,
    RENDERER_ADDED_SLIDES,
)
from app.shared.models import as_utc

logger = logging.getLogger(__name__)

# Пределы длины НА РЕНДЕРЕ. Схема ответа модели свои границы уже проверила;
# эти относятся к тому, что в схеме не описано, — именам документов, имени
# блокнота, заголовку плана.
TITLE_MAX_CHARS = 120
SUBTITLE_MAX_CHARS = 200
HEADING_MAX_CHARS = 120
BULLET_MAX_CHARS = 300
SOURCE_LINE_MAX_CHARS = 200
# Слайд «Источники» не умеет прокручиваться: строки сверх этого числа не
# поместятся на нём физически, поэтому последняя строка честно говорит,
# сколько документов не показано, вместо молчаливого обрезания списка.
MAX_SOURCE_LINES = 18

# Раскладки, которыми пайплайн пользуется, когда реестра шаблонов ещё нет
# (модуль templates.py делается отдельно) или в нём нет запрошенного ключа.
# Индексы — из стандартной темы python-pptx: 0 «Title Slide», 1 «Title and
# Content», 2 «Section Header».
_FALLBACK_LAYOUTS = {"title": 0, "section": 2, "bullets": 1, "sources": 1}


@dataclass
class RenderedSource:
    """Документ в списке источников: имя и страницы, на которые ссылались."""

    source_id: int
    name: str
    pages: list[int] = field(default_factory=list)


def fit(value: str, limit: int) -> str:
    """Строка не длиннее limit. Обрезка, а не отказ.

    Уронить готовую колоду из-за длинного имени файла — худший из возможных
    исходов: работа модели уже сделана и оплачена временем пользователя.
    """
    text = " ".join((value or "").split())
    if len(text) <= limit:
        return text
    return text[: limit - 1].rstrip() + "…"


def resolve_template(template_key: str) -> Any | None:
    """TemplateInfo из реестра шаблонов или None.

    Импорт отложенный и защищённый намеренно. Реестр (templates.py) делается
    параллельно, и до его появления пайплайн обязан оставаться рабочим:
    отсутствие оформления — это повод нарисовать колоду темой по умолчанию, а
    не повод не отдать пользователю ничего. Как только модуль появится, этот
    же код начнёт брать шаблон из него без единой правки.
    """
    try:
        from app.modules.presentations.templates import template_registry
    except ImportError:
        logger.warning(
            "Template registry is not available yet, rendering with the default "
            "python-pptx theme (template_key=%r)",
            template_key,
        )
        return None
    template = template_registry.get(template_key)
    if template is None:
        logger.warning(
            "Unknown presentation template %r, rendering with the default theme",
            template_key,
        )
    return template


def _new_presentation(template: Any | None):
    from pptx import Presentation

    if template is None:
        return Presentation(), dict(_FALLBACK_LAYOUTS)
    presentation = Presentation(str(template.template_file))
    layouts = dict(_FALLBACK_LAYOUTS)
    # Ключи реестра перекрывают умолчания по одному: шаблон, у которого не
    # объявлена, скажем, раскладка section, всё равно рисуется — просто на
    # раскладке по умолчанию.
    layouts.update(getattr(template, "layouts", None) or {})
    return presentation, layouts


def _layout(presentation, layouts: dict[str, int], name: str):
    """Раскладка по имени из реестра, с откатом на первую доступную.

    Индекс из реестра может не совпасть с содержимым файла (шаблон
    перерисовали, раскладку удалили). Это повод нарисовать слайд не тем
    макетом, а не уронить генерацию.
    """
    available = presentation.slide_layouts
    index = layouts.get(name, _FALLBACK_LAYOUTS.get(name, 0))
    if not isinstance(index, int) or index < 0 or index >= len(available):
        logger.warning(
            "Layout %r resolves to index %r, which the template does not have; "
            "falling back to layout 0",
            name,
            index,
        )
        index = 0
    return available[index]


def _set_title(slide, text: str) -> None:
    title = getattr(slide.shapes, "title", None)
    if title is None:
        # Раскладка без заголовка — редкость, но встречается в чужих шаблонах.
        # Заголовок важнее оформления, поэтому кладём его в тело.
        frame = _body_frame(slide)
        if frame is not None:
            frame.paragraphs[0].text = text
        return
    title.text = text


def _is_title(placeholder, title) -> bool:
    """Тот ли это плейсхолдер, что и заголовок слайда.

    Сравнение по shape_id, а НЕ по `is`: python-pptx создаёт новый объект-обёртку
    на каждое обращение к фигуре, поэтому `placeholder is slide.shapes.title`
    ложно даже для одного и того же элемента. С таким сравнением заголовок не
    исключался из перебора, и текст тела затирал его — на титульном слайде
    вместо названия презентации оставалась подпись.
    """
    if title is None:
        return False
    try:
        return placeholder.shape_id == title.shape_id
    except AttributeError:  # pragma: no cover - защита от чужой реализации
        return placeholder is title


def _body_frame(slide):
    """Текстовое поле слайда: первый плейсхолдер, который не заголовок.

    По индексу 1 обращаться нельзя: в чужих шаблонах у раскладки бывает другой
    набор плейсхолдеров, и обращение по индексу роняет рендер KeyError'ом.
    """
    title = getattr(slide.shapes, "title", None)
    for placeholder in slide.placeholders:
        if _is_title(placeholder, title):
            continue
        if placeholder.has_text_frame:
            return placeholder.text_frame
    # Плейсхолдера нет вовсе — рисуем собственную рамку. Оформление хуже, чем
    # у шаблона, но текст слайда доходит до пользователя. Размеры в дюймах, а
    # не в долях слайда: до размеров слайда из объекта слайда надо идти через
    # package, и эта дорога зависит от версии python-pptx сильнее, чем сама
    # рамка того стоит.
    from pptx.util import Inches

    box = slide.shapes.add_textbox(Inches(0.8), Inches(1.8), Inches(8.4), Inches(4.5))
    return box.text_frame


def _fill_bullets(slide, bullets: list[str], limit: int) -> None:
    frame = _body_frame(slide)
    if frame is None:  # pragma: no cover - _body_frame всегда что-то возвращает
        return
    frame.clear()
    for index, bullet in enumerate(bullets or ["—"]):
        paragraph = frame.paragraphs[0] if index == 0 else frame.add_paragraph()
        paragraph.text = fit(bullet, limit)
        paragraph.level = 0


def format_source_line(source: RenderedSource, language: str) -> str:
    """Строка списка источников: «Имя документа — стр. 3, 5»."""
    name = fit(source.name or f"#{source.source_id}", SOURCE_LINE_MAX_CHARS - 40)
    pages = [page for page in sorted(set(source.pages)) if isinstance(page, int)]
    if not pages:
        return name
    label = PAGES_LABEL.get(language, PAGES_LABEL["ru"])
    return fit(f"{name} — {label} {', '.join(str(page) for page in pages)}", SOURCE_LINE_MAX_CHARS)


def build_source_lines(sources: list[RenderedSource], language: str) -> list[str]:
    lines = [format_source_line(source, language) for source in sources]
    if len(lines) <= MAX_SOURCE_LINES:
        return lines or ["—"]
    hidden = len(lines) - MAX_SOURCE_LINES
    return [*lines[:MAX_SOURCE_LINES], f"… +{hidden}"]


def render_presentation(
    *,
    title: str,
    slides: list[PresentationSlide],
    sources: list[RenderedSource],
    language: str,
    template_key: str,
    notebook_name: str,
    created_at: datetime,
    output_path: str,
) -> None:
    """Собрать колоду и сохранить её по output_path.

    Состав файла фиксирован контрактом slide_count: титульный слайд, по слайду
    на каждую секцию плана и финальные «Источники» — ровно
    len(slides) + RENDERER_ADDED_SLIDES слайдов.

    Раскладка section из реестра здесь не используется: разделители увеличили
    бы число слайдов сверх заказанного пользователем. Ключ остаётся в
    контракте реестра для колод с разделами, если такие появятся.
    """
    template = resolve_template(template_key)
    presentation, layouts = _new_presentation(template)

    # --- Титульный слайд ---
    title_slide = presentation.slides.add_slide(_layout(presentation, layouts, "title"))
    _set_title(title_slide, fit(title, TITLE_MAX_CHARS))
    slide_total = len(slides) + RENDERER_ADDED_SLIDES
    subtitle = " · ".join(
        part
        for part in (
            fit(notebook_name, 100),
            as_utc(created_at).strftime("%d.%m.%Y"),
            f"{slide_total} {SLIDES_WORD.get(language, SLIDES_WORD['ru'])}",
        )
        if part
    )
    _fill_subtitle(title_slide, fit(subtitle, SUBTITLE_MAX_CHARS))

    # --- Контентные слайды ---
    bullets_layout = _layout(presentation, layouts, "bullets")
    for slide_data in slides:
        slide = presentation.slides.add_slide(bullets_layout)
        _set_title(slide, fit(slide_data.heading, HEADING_MAX_CHARS))
        _fill_bullets(slide, list(slide_data.bullets), BULLET_MAX_CHARS)

    # --- Источники ---
    sources_slide = presentation.slides.add_slide(
        _layout(presentation, layouts, "sources")
    )
    _set_title(sources_slide, SOURCES_HEADING.get(language, SOURCES_HEADING["ru"]))
    _fill_bullets(
        sources_slide, build_source_lines(sources, language), SOURCE_LINE_MAX_CHARS
    )

    presentation.save(output_path)


def _fill_subtitle(slide, text: str) -> None:
    """Подпись титульного слайда, если для неё есть место.

    Отдельной веткой от _fill_bullets: на титульной раскладке подпись — это
    один абзац, и добавлять к нему пустую рамку, когда плейсхолдера нет, не
    нужно. Титул без подписи — рабочий слайд.
    """
    title = getattr(slide.shapes, "title", None)
    for placeholder in slide.placeholders:
        if _is_title(placeholder, title):
            continue
        if placeholder.has_text_frame:
            placeholder.text_frame.text = text
            return
